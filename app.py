import streamlit as st
import json
import io
import os
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt

# --- 1. 日本語フォントの設定 ---
def setup_japanese_font():
    # フォントファイルのパス（リポジトリ内の相対パス）
    font_path = "fonts/ipaexg.ttf" 
    
    if os.path.exists(font_path):
        # フォントマネージャーに追加
        fm.fontManager.addfont(font_path)
        # フォントプロパティを取得してMatplotlibのデフォルトに設定
        font_prop = fm.FontProperties(fname=font_path)
        plt.rcParams['font.family'] = font_prop.get_name()
        return font_prop
    else:
        st.warning("日本語フォント(fonts/ipaexg.ttf)が見つかりません。文字化けする可能性があります。")
        return None

# フォント設定を実行
jp_font = setup_japanese_font()

# --- 2. スライド生成ロジック ---
def create_slide_deck(json_data):
    prs = Presentation()
    
    for slide_data in json_data:
        # 白紙スライド (Layout 6)
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = slide_data.get('title', 'No Title')
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True

        sType = slide_data.get('type')
        content = slide_data.get('content', {})
        
        # --- Type A: 箇条書き (Standard) ---
        if sType == 'bullet_points':
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for item in content.get('points', []):
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.space_after = Pt(10)

        # --- Type B: 棒グラフ (Simple Chart) ---
        elif sType == 'bar_chart':
            fig, ax = plt.subplots(figsize=(8, 4.5))
            labels = content.get('labels', [])
            values = content.get('values', [])
            
            # デザイン調整
            ax.bar(labels, values, color='#4A90E2', alpha=0.8)
            ax.set_title(slide_data.get('title'), fontsize=14)
            ax.grid(axis='y', linestyle='--', alpha=0.5)
            
            # 画像として保存してスライドへ
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=150)
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))
            plt.close()

       # --- Type C: ネットワーク図 (修正版: 頑丈なデータ処理) ---
        elif sType == 'network_graph':
            fig, ax = plt.subplots(figsize=(8, 5))
            G = nx.Graph()
            
            # --- データの正規化処理 ---
            raw_nodes = content.get('nodes', [])
            raw_edges = content.get('edges', [])
            
            clean_nodes = []
            clean_edges = []

            # Nodesのクリーニング（文字列のリストにする）
            for n in raw_nodes:
                if isinstance(n, str):
                    clean_nodes.append(n)
                elif isinstance(n, list) and len(n) > 0:
                    clean_nodes.append(str(n[0])) # ["A"] -> "A"
                elif isinstance(n, dict):
                    # {"name": "A"} -> "A" (最初の値を採用)
                    clean_nodes.append(str(list(n.values())[0]))

            # Edgesのクリーニング（[source, target]のリストにする）
            for e in raw_edges:
                if isinstance(e, list) and len(e) >= 2:
                    clean_edges.append((str(e[0]), str(e[1])))
                elif isinstance(e, dict):
                    # {"source": "A", "target": "B"} -> ("A", "B")
                    vals = list(e.values())
                    if len(vals) >= 2:
                        clean_edges.append((str(vals[0]), str(vals[1])))

            # グラフの構築（データがあれば）
            if clean_nodes:
                G.add_nodes_from(clean_nodes)
            if clean_edges:
                G.add_edges_from(clean_edges)
            
            # データが空の場合の安全策
            if G.number_of_nodes() == 0:
                ax.text(0.5, 0.5, "No Data for Graph", ha='center', va='center')
            else:
                # レイアウト計算
                try:
                    pos = nx.spring_layout(G, k=0.8, seed=42)
                    nx.draw_networkx_nodes(G, pos, node_size=2000, node_color='#E8F5E9', edgecolors='#2E7D32', ax=ax)
                    nx.draw_networkx_edges(G, pos, width=2, edge_color='#90A4AE', ax=ax)
                    # フォント設定の安全策
                    f_family = jp_font.get_name() if jp_font else 'sans-serif'
                    nx.draw_networkx_labels(G, pos, font_family=f_family, font_size=11, ax=ax)
                except Exception as ex:
                    # 描画エラー時も落ちないようにする
                    ax.text(0.5, 0.5, f"Graph Error: {ex}", ha='center', va='center')
            
            ax.axis('off')
            ax.set_title("Concept Map", fontsize=14, loc='left', color='gray')
            
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=150)
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))
            plt.close()

       # --- Type D: タイムライン (修正版: 頑丈なデータ処理) ---
        elif sType == 'timeline':
            fig, ax = plt.subplots(figsize=(8, 3))
            
            # コンテンツの取得（eventsキーがない場合はcontentそのものをリストとして扱う）
            events = content.get('events', [])
            if not events and isinstance(content, list):
                events = content
            
            dates = []
            labels = []

            # 柔軟なキー読み取り処理
            for e in events:
                if isinstance(e, dict):
                    # 日付っぽいキーを探す
                    d = e.get('date') or e.get('year') or e.get('time') or e.get('Date') or "N/A"
                    # ラベルっぽいキーを探す
                    l = e.get('label') or e.get('title') or e.get('event') or e.get('Label') or "No Label"
                    
                    dates.append(str(d))
                    labels.append(str(l))
            
            # データがある場合のみ描画
            if dates:
                # 簡易的なタイムライン描画
                ax.hlines(1, 0, len(dates)-1, color='#FF7043', linewidth=3) # メインライン
                ax.plot(range(len(dates)), [1]*len(dates), 'o', markersize=10, color='#FF7043') # 点
                
                # テキスト配置
                for i, (date, label) in enumerate(zip(dates, labels)):
                    # 日付（上側）
                    ax.text(i, 1.15, date, ha='center', fontsize=10, color='gray')
                    # ラベル（下側・折り返し対応）
                    # 長いラベルは改行を入れる簡易処理
                    display_label = label[:10] + '...' if len(label) > 10 else label
                    ax.text(i, 0.85, display_label, ha='center', va='top', fontsize=12, fontweight='bold')
                    
                ax.axis('off')
                ax.set_ylim(0.5, 1.5)
            else:
                ax.text(0.5, 0.5, "No Timeline Data", ha='center', va='center')
                ax.axis('off')
            
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', bbox_inches='tight', dpi=150)
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(2.5), width=Inches(8))
            plt.close()

    # PPTX保存処理
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- 3. Streamlit UI ---
st.set_page_config(page_title="AI Visual Slide Gen", layout="wide")

col1, col2 = st.columns([1, 1])

with col1:
    st.title("🤖 Info-Graph Slide Generator")
    st.markdown("""
    **Geminiで作ったJSONを貼るだけ**で、構造化されたビジュアルスライドを生成します。
    NotebookLMのような概念図や、タイムラインも自動描画します。
    """)
    
    st.info("💡 使い方: Geminiに専用プロンプトを投げて、出てきたJSONを右のボックスに貼り付けてください。")

with col2:
    json_input = st.text_area("JSON Input:", height=400, placeholder='[Paste JSON code here...]')
    
    if st.button("🚀 Generate Slides", type="primary"):
        if json_input:
            with st.spinner('Generating visuals & slides...'):
                try:
                    data = json.loads(json_input)
                    pptx_file = create_slide_deck(data)
                    
                    st.success("完了しました！")
                    st.download_button(
                        label="📥 Download .pptx",
                        data=pptx_file,
                        file_name="visual_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except json.JSONDecodeError:
                    st.error("JSONの形式が間違っています。括弧の閉じ忘れなどを確認してください。")
                except Exception as e:
                    st.error(f"エラーが発生しました: {e}")
